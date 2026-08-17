"""Build the field-day results slide deck from already-curated, already-verified material.

    python make_slide_deck.py Data_NatureSpec/2026_Aug_16 --out KOTZEBUE_20260816_RESULTS.pptx

Every number and figure here already exists and was already checked elsewhere this
session -- this script does no new analysis. It reads from `highlights/` (make_highlights.py,
run first) for the images, and its bullet text is transcribed from GLOBAL_COMPARISON.md,
PAPER_READINESS.md and the per-station *_GIOP_FINDINGS.md, not invented. Run
`make_highlights.py` on the same day folder before this.
"""

import argparse
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W, H = Inches(13.333), Inches(7.5)   # 16:9
NAVY = RGBColor(0x1a, 0x2e, 0x4a)
TEAL = RGBColor(0x1f, 0x7a, 0x99)
GREY = RGBColor(0x55, 0x55, 0x55)
RED = RGBColor(0xc0, 0x39, 0x2b)
WHITE = RGBColor(0xff, 0xff, 0xff)


def _title(slide, text, sub=None):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(15); r2.font.color.rgb = GREY
    return box


def _bullets(slide, items, left, top, width, height, size=14):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, level, color, bold = item
        else:
            text, level, color, bold = item, 0, RGBColor(0x22, 0x22, 0x22), False
        p.level = level
        r = p.add_run(); r.text = ("• " if level == 0 else "– ") + text
        r.font.size = Pt(size - level)
        r.font.color.rgb = color
        r.font.bold = bold
        p.space_after = Pt(6)
    return box


def _pic_fit(slide, path, left, top, max_w, max_h):
    """Place an image scaled to fit inside a box, centred, preserving aspect ratio."""
    if not os.path.exists(path):
        box = slide.shapes.add_textbox(left, top, max_w, max_h)
        p = box.text_frame.paragraphs[0]
        r = p.add_run(); r.text = "[missing: %s]" % os.path.basename(path)
        r.font.color.rgb = RED; r.font.size = Pt(12)
        return None
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w; h = Emu(int(max_w / ar))
    else:
        h = max_h; w = Emu(int(max_h * ar))
    x = left + Emu(int((max_w - w) / 2))
    y = top + Emu(int((max_h - h) / 2))
    return slide.shapes.add_picture(path, x, y, width=w, height=h)


def _caption(slide, text, left, top, width, size=11, color=GREY, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.italic = True; r.font.color.rgb = color


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("day_folder")
    ap.add_argument("--out", default="RESULTS.pptx")
    a = ap.parse_args()
    day = a.day_folder.rstrip("/")
    hl = os.path.join(day, "highlights")
    byloc = os.path.join(day, "by_location")

    def station_hl(station, fn):
        return os.path.join(byloc, station, "analysis", "highlights", fn)

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # ---------------------------------------------------------------- 1. Title
    s = prs.slides.add_slide(blank)
    box = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2.2))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Above-Water Remote-Sensing Reflectance & GIOP Inversion"
    r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = "Kotzebue, Alaska  ·  2026-08-16  ·  3 sites, 6 validated R_rs products"
    r2.font.size = Pt(20); r2.font.color.rgb = TEAL
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = ("NaturaSpec Plus field spectroradiometry, angle-matched sky pairing, "
              "GIOP semi-analytical inversion (tmittal22/giop-workbench)")
    r3.font.size = Pt(14); r3.font.color.rgb = GREY

    # ---------------------------------------------------------------- 2. Map
    s = prs.slides.add_slide(blank)
    _title(s, "Where: three sites, two foreoptics",
          "60 scans total  ·  LOC1 (22), LOC2 (23, split a/b/c), LOC3 (15, FIBR15 15° + FLENS8 8°, split main/murky)")
    _pic_fit(s, os.path.join(hl, "site_map.png"), Inches(0.4), Inches(1.15),
            Inches(12.5), Inches(6.0))

    # ---------------------------------------------------------------- Method: rho / angle
    LOC1 = "LOC1_66.89718N_162.60290W/FLENS8_FOV08"
    s = prs.slides.add_slide(blank)
    _title(s, "Method — the ρ angle correction",
          "LOC1 example: rho_at_angle scales Mobley (1999) ρ=0.028 by the Fresnel ratio at the scan's own view angle")
    _pic_fit(s, os.path.join(byloc, LOC1, "analysis", "fig8_rho_angle_correction.png"),
            Inches(0.6), Inches(1.25), Inches(12.0), Inches(4.7))
    _bullets(s, [
        "42–50° is an ordinary hand-aimed spread; Fresnel reflectance rises ~32% across it — a fixed ρ puts all of that into R_rs as a pointing-tied systematic",
        "Correcting each LOC1 scan at its own angle removed a p=0.006 tilt-vs-R_rs trend and cut scatter 11.1% → 8.6%",
        "Sky-side angle-matching, by contrast, did NOT matter at any of the 7 stations — every station's sky scans spanned only 3–6°, never enough range to bite on (SKY_CHOICE_SYNTHESIS.md)",
    ], Inches(0.4), Inches(5.9), Inches(12.5), Inches(1.5), size=15)

    # ---------------------------------------------------------------- Method: scaled mean
    s = prs.slides.add_slide(blank)
    _title(s, "Method — shape-consistency: what does \"the mean\" mean?",
          "LOC1 example: the amplitude-normalised (scaled) mean, not the plain mean")
    _pic_fit(s, os.path.join(byloc, LOC1, "analysis", "fig11_scaled_mean_method.png"),
            Inches(0.6), Inches(1.2), Inches(12.0), Inches(4.3))
    _bullets(s, [
        "Replicate scans differ in both SHAPE (what the water is) and AMPLITUDE (how bright it happened to be, e.g. sun flicker) — a plain mean conflates the two and misreports uncertainty",
        "The scaled-mean estimator splits them: R_rs(λ) = S(λ) × c, with S (shape) and c (amplitude) fitted separately, iterated to convergence",
        "At LOC1: shape spread ±1.7% (\"what the water is\") vs amplitude spread ±11% (\"how much\") — very different numbers a plain mean would blur into one",
    ], Inches(0.4), Inches(5.6), Inches(12.5), Inches(1.8), size=14)

    # ---------------------------------------------------------------- Method: GIOP good vs bad
    s = prs.slides.add_slide(blank)
    _title(s, "GIOP fits: why LOC1 is the worst fit and LOC2 the best")
    _pic_fit(s, station_hl(LOC1, "giop6_all_fits.png"), Inches(0.3), Inches(1.2),
            Inches(6.3), Inches(3.3))
    _caption(s, "LOC1 — worst fit of any station", Inches(0.3), Inches(4.55), Inches(6.3))
    _pic_fit(s, station_hl("LOC2a_66.89677N_162.57953W/FLENS8_FOV08", "giop6_all_fits.png"),
            Inches(6.7), Inches(1.2), Inches(6.3), Inches(3.3))
    _caption(s, "LOC2a — best fit of any station", Inches(6.7), Inches(4.55), Inches(6.3))
    _bullets(s, [
        ("Does releasing extra freedom (fitted a*_φ family/seed) actually help? It depends on the station:", 0, NAVY, True),
        ("LOC1: barely. χ²_ν 14.0 (free) → 13.9 (max freedom), RMS misfit 8.5% → 8.6% — essentially flat. Almost all of LOC1's gain happens earlier, constrained→free (χ²_ν 48.0→14.0).", 1, RGBColor(0x22,0x22,0x22), False),
        ("LOC2a: substantially. χ²_ν 104.3 (free) → 23.6 (max freedom) — a real 4.4× drop, RMS 4.7%→2.9%.", 1, RGBColor(0x22,0x22,0x22), False),
        ("LOC2b: same pattern. χ²_ν 6199 → 999, a 6.2× drop, RMS 4.2%→2.9%.", 1, RGBColor(0x22,0x22,0x22), False),
        ("LOC1's residual is NOT a missing-parameter problem — it's traced to the blue band specifically (12.75% RMS there alone, worst +38%), where LOC1's high a_dg makes a wrong CDOM-slope shape assumption costly", 0, RED, True),
        ("Never compare χ²_ν itself across stations (it scales with each station's own measured σ) — RMS misfit is the number above that is actually comparable", 0, TEAL, False),
    ], Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.4), size=14)

    # ---------------------------------------------------------------- 3/4/5. Per-site
    site_defs = [
        ("LOC1_66.89718N_162.60290W/FLENS8_FOV08", "LOC1",
         "Reference station  ·  22 scans (12 water, 8 sky, 2 land targets)  ·  FLENS8 (8° FOV)",
         [
             ("R_rs(555) = 0.00863 sr⁻¹  ·  one clean water population (12/12 kept)", 0, RGBColor(0x22,0x22,0x22), False),
             ("GIOP (free): a_dg(443)=0.78 m⁻¹, b_bp(443)=0.043 m⁻¹  ·  highest CDOM of any station", 0, RGBColor(0x22,0x22,0x22), False),
             ("Best-fit RMS misfit 8.5% — worst of the 6 stations, traced to the blue band (12.75% RMS there alone)", 0, RGBColor(0x22,0x22,0x22), False),
             ("00005/00007 were glint-correctable (same test/fix as LOC2a's 00035) — applied: shape consistency 7.08%→3.68%, R_rs(555) moved only −0.23%, GIOP composition unchanged", 0, TEAL, False),
         ]),
        ("LOC2a_66.89677N_162.57953W/FLENS8_FOV08", "LOC2a — main open water",
         "9 scans  ·  glint-corrected (scan 00035, Ruddick et al. 2006 NIR-similarity)",
         [
             ("R_rs(555) = 0.01066 sr⁻¹", 0, RGBColor(0x22,0x22,0x22), False),
             ("GIOP (free): a_dg(443)=0.49 m⁻¹, b_bp(443)=0.041 m⁻¹", 0, RGBColor(0x22,0x22,0x22), False),
             ("Best RMS misfit 2.9% — much better than LOC1, reproduces per-scan (9/9 improve)", 0, RGBColor(0x22,0x22,0x22), False),
         ]),
    ]
    for station, name, sub, bullets in site_defs:
        s = prs.slides.add_slide(blank)
        _title(s, name, sub)
        _pic_fit(s, station_hl(station, "fig12_FINAL_mean_Rrs.png"),
                Inches(0.3), Inches(1.2), Inches(5.9), Inches(3.3))
        _caption(s, "Final R_rs (scaled-mean method)", Inches(0.3), Inches(4.55), Inches(5.9))
        _pic_fit(s, station_hl(station, "giop10_final_result.png"),
                Inches(6.3), Inches(1.2), Inches(6.6), Inches(3.3))
        _caption(s, "GIOP inversion, final result", Inches(6.3), Inches(4.55), Inches(6.6))
        _bullets(s, bullets, Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.2), size=15)

    # LOC2b + LOC2c combined slide (disturbed water + algae, the "not open water" cases)
    s = prs.slides.add_slide(blank)
    _title(s, "LOC2b & LOC2c — not the same water as LOC2a",
          "Both split off LOC2's original 23-scan population by analyse_water_scans.py's QC test")
    _pic_fit(s, station_hl("LOC2b_66.89677N_162.57953W/FLENS8_FOV08", "fig12_FINAL_mean_Rrs.png"),
            Inches(0.3), Inches(1.15), Inches(5.9), Inches(2.9))
    _caption(s, "LOC2b: disturbed/turbid water (n=3, R_rs)", Inches(0.3), Inches(4.05), Inches(5.9))
    _pic_fit(s, os.path.join(hl, "field_photos", "LOC2c_00043.jpg"),
            Inches(6.3), Inches(1.15), Inches(3.2), Inches(2.9))
    _caption(s, "LOC2c: floating algae mat (photo)", Inches(6.3), Inches(4.05), Inches(3.2))
    _pic_fit(s, os.path.join(byloc, "LOC2c_66.89677N_162.57953W/FLENS8_FOV08/analysis/c1_algae_reflectance.png"),
            Inches(9.6), Inches(1.15), Inches(3.3), Inches(2.9))
    _caption(s, "LOC2c: reflectance, NOT R_rs (opaque target)", Inches(9.6), Inches(4.05), Inches(3.3))
    _bullets(s, [
        "LOC2b (n=3): recently disturbed sediment, confirmed by timing (first 3 scans, settling over ~2 min) — real water, not an artefact",
        "b_bp(443)=0.047 m⁻¹ (vs LOC2a's 0.041) — elevated backscatter, corroborated 3 independent ways: spectral clustering, a failed glint-collapse test, and now GIOP composition",
        "LOC2c (n=2): a floating algae clump is opaque enough that subtracting ρ·L_sky would remove real signal, not a contaminant — reflectance R = π·L_target/E_d is the correct product, not R_rs",
    ], Inches(0.4), Inches(4.5), Inches(12.5), Inches(2.6), size=15)

    # LOC3 site slide
    s = prs.slides.add_slide(blank)
    _title(s, "LOC3 — three water populations, one open caveat",
          "FIBR15 (15° FOV): main (n=3) + murky (n=2, 00058/00059)  ·  FLENS8 (8° FOV): n=4")
    _pic_fit(s, os.path.join(hl, "LOC3_site_summary.png"), Inches(0.3), Inches(1.15),
            Inches(12.7), Inches(3.5))
    _bullets(s, [
        ("GIOP (free), a_dg(443) / b_bp(443) / best RMS:", 0, NAVY, True),
        ("FIBR15 main:  1.10 m⁻¹  /  0.079 m⁻¹  /  4.9%", 1, RGBColor(0x22,0x22,0x22), False),
        ("FIBR15 murky:  1.36 m⁻¹  /  0.122 m⁻¹  /  2.8%  — highest b_bp of any station, 2nd independent confirmation of the sediment split", 1, RGBColor(0x22,0x22,0x22), False),
        ("FLENS8:  1.22 m⁻¹  /  0.091 m⁻¹  /  5.2%", 1, RGBColor(0x22,0x22,0x22), False),
        ("⚠ Every LOC3 sub-population shows a 2nd R_rs peak near 805–810 nm — optically deep water cannot produce this", 0, RED, True),
        ("Confirmed in raw radiance (not a processing artefact): L(810)/L(700) = 0.416 at LOC3 vs 0.226 at LOC1", 1, RED, False),
        ("GIOP only fits 400–700 nm, so a \"good\" GIOP fit here is NOT evidence against the bottom caveat — it simply never sees the 810 nm evidence", 1, RED, False),
    ], Inches(0.4), Inches(4.75), Inches(12.5), Inches(2.6), size=14)

    # LOC3 bottom-visibility photo slide
    s = prs.slides.add_slide(blank)
    _title(s, "What the LOC3 photos show", "Same field day, same instrument — visibly different water")
    _pic_fit(s, os.path.join(hl, "field_photos", "LOC3_00055.jpg"),
            Inches(0.4), Inches(1.3), Inches(4.0), Inches(3.6))
    _caption(s, "FIBR15 main — open water", Inches(0.4), Inches(5.0), Inches(4.0))
    _pic_fit(s, os.path.join(hl, "field_photos", "LOC3_00058.jpg"),
            Inches(4.6), Inches(1.3), Inches(4.0), Inches(3.6))
    _caption(s, "FIBR15 murky — the split is spectral, not always obvious by eye", Inches(4.6), Inches(5.0), Inches(4.0))
    _pic_fit(s, os.path.join(hl, "field_photos", "LOC3_00047.jpg"),
            Inches(8.8), Inches(1.3), Inches(4.0), Inches(3.6))
    _caption(s, "FLENS8 — substrate visible through the water column", Inches(8.8), Inches(5.0), Inches(4.0))

    # ---------------------------------------------------------------- Global analysis
    s = prs.slides.add_slide(blank)
    _title(s, "Global comparison — R_rs across all 6 stations")
    _pic_fit(s, os.path.join(hl, "all_stations_overplot.png"), Inches(0.4), Inches(1.15),
            Inches(12.5), Inches(4.6))
    _bullets(s, [
        "LOC1 and LOC2 (both sub-populations) collapse smoothly toward zero past ~700 nm, as optically deep water must",
        "LOC3 does not — a structurally different measurement regime, not a point on the same continuum",
    ], Inches(0.4), Inches(5.9), Inches(12.5), Inches(1.3), size=15)

    s = prs.slides.add_slide(blank)
    _title(s, "Global comparison — GIOP composition across all 6 stations")
    _pic_fit(s, os.path.join(hl, "giop_cross_station_all.png"), Inches(0.4), Inches(1.15),
            Inches(12.5), Inches(4.6))
    _bullets(s, [
        "LOC1/LOC2: three real, distinct water states, each independently corroborated (not GIOP artefacts)",
        "LOC3: a_dg and b_bp are elevated at ALL THREE sub-populations, including the clean ones — more consistent with GIOP absorbing a bright shallow bottom into its only available water-column parameters than with a real compositional signal",
    ], Inches(0.4), Inches(5.9), Inches(12.5), Inches(1.3), size=15)

    s = prs.slides.add_slide(blank)
    _title(s, "What generalizes beyond this one field day")
    _bullets(s, [
        "χ²_ν is NOT comparable across stations with different measured shape uncertainty — a χ²_ν of 16,000 at LOC2b vs 74 at LOC1 does NOT mean LOC2b fits 200× worse (its RMS misfit is actually better, 9.7% vs 10.9%). Compare RMS misfit instead.",
        "A profile-solver bug can look exactly like \"the data can't determine the shapes\" — always check the nesting property (free ≥ as good as constrained) before trusting that conclusion.",
        "A raw-radiance ratio, L(NIR-far)/L(NIR-near), is a cheap, processing-independent screen for bottom contamination in any above-water R_rs dataset.",
        "The sky-scan angle choice did not matter at any of the 7 stations this field day — every station's sky scans spanned only 3–6°, never enough range for angle-matching to have anything to bite on. (The water-side angle correction is a separate, validated result and is NOT affected — it fixed a real 27% effect at LOC1.)",
        "Wind speed was never recorded live, but was recovered after the fact from the Kotzebue airport's public ASOS archive for every station — a free, retroactive fix worth trying at any past coastal site.",
    ], Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.6), size=17)

    s = prs.slides.add_slide(blank)
    _title(s, "Open items, ranked by value per unit of future field effort")
    _bullets(s, [
        ("1. A depth measurement at LOC3", 0, NAVY, True),
        ("Cheapest fix available (a weighted line) — the only thing that converts the bottom-reflectance caveat into a resolved question.", 1, RGBColor(0x22,0x22,0x22), False),
        ("2. Record wind speed live at any future site not near an ASOS station", 0, NAVY, True),
        ("3. A same-instant, same-range, two-foreoptic comparison", 0, NAVY, True),
        ("To finally answer the footprint/FOV question this field day could not (range, time, and bottom depth were all confounded at once).", 1, RGBColor(0x22,0x22,0x22), False),
        ("4. Refine LOC2a's fit_shapes grid", 0, NAVY, True),
        ("The solver lands ~5% off the true χ² minimum there — a small, real implementation gap.", 1, RGBColor(0x22,0x22,0x22), False),
    ], Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.9), size=16)

    n_slides = len(prs.slides._sldIdLst)
    prs.save(a.out)
    print("wrote", a.out, "(%d slides)" % n_slides)


if __name__ == "__main__":
    main()
